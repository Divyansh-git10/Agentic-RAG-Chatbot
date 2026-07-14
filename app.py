import streamlit as st
import os
import tempfile

from agents import IngestionAgent, RetrievalAgent, LLMResponseAgent
from mcp import create_mcp_message


def extract_text_from_upload(file):
    """Extract plain text from an uploaded file based on its extension.

    Supports .pdf, .docx, .pptx, .csv, .txt and .md, falling back to a
    plain-text decode for anything else. Parsing libraries are imported
    lazily so the app still starts if an optional dependency is missing,
    and a failure on one file never breaks the rest of the batch.
    """
    import io

    name = (file.name or "").lower()
    data = file.read()

    try:
        if name.endswith(".pdf"):
            import fitz  # PyMuPDF
            doc = fitz.open(stream=data, filetype="pdf")
            return "\n".join(page.get_text() for page in doc)

        if name.endswith(".docx"):
            from docx import Document
            document = Document(io.BytesIO(data))
            return "\n".join(p.text for p in document.paragraphs)

        if name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(data))
            return "\n".join(
                shape.text_frame.text
                for slide in prs.slides
                for shape in slide.shapes
                if shape.has_text_frame
            )

        if name.endswith(".csv"):
            import csv
            text_data = data.decode("utf-8", errors="replace")
            return "\n".join(", ".join(row) for row in csv.reader(io.StringIO(text_data)))

        # .txt, .md, and anything else: decode as text
        try:
            return data.decode("utf-8")
        except UnicodeDecodeError:
            return data.decode("latin1")
    except Exception as exc:  # one bad file shouldn't break the whole upload
        return f"[Could not parse {file.name}: {exc}]"


st.set_page_config(page_title="RAG Chatbot", layout="wide")
st.title("📄 Agentic RAG Chatbot")


# --- File Upload ---
st.sidebar.header("📤 Upload Documents")
uploaded_files = st.sidebar.file_uploader("Upload multiple files", accept_multiple_files=True)
docs = []
for file in uploaded_files:
    docs.append(extract_text_from_upload(file))


# Initialize session state
if "retriever" not in st.session_state:
    st.session_state.retriever = None
if "responder" not in st.session_state:
    st.session_state.responder = LLMResponseAgent()
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []

# On file upload
if uploaded_files:
    ingestor = IngestionAgent()
    trace_id = "streamlit-chat"
    msg_ingested = ingestor.handle({
        "trace_id": trace_id,
        "payload": {
            "docs": docs  
        }
    })

    retriever = RetrievalAgent()
    retriever.handle(msg_ingested)

    st.session_state.retriever = retriever
    st.success("✅ Documents parsed and indexed!")


# --- Chat Interface ---
st.subheader("💬 Ask a Question")

query = st.text_input("Enter your question:")

if st.button("Ask") and query and st.session_state.retriever:
    trace_id = "streamlit-chat"
    query_msg = create_mcp_message("UI", "RetrievalAgent", "QUERY", trace_id, {"query": query})
    retrieval_response = st.session_state.retriever.handle(query_msg)

    final_answer = st.session_state.responder.handle(retrieval_response)

    # Show result
    st.markdown(f"### 🧠 Answer:\n{final_answer['answer']}")

    # Source Chunks
    with st.expander("📚 Source Context"):
        for i, chunk in enumerate(final_answer["source_chunks"]):
            st.markdown(f"**Chunk {i+1}:** {chunk.strip()}\n")

    # Save chat
    st.session_state.chat_history.append((query, final_answer["answer"]))

# Chat history
if st.session_state.chat_history:
    st.sidebar.markdown("### 🕒 Chat History")
    for q, a in reversed(st.session_state.chat_history):
        st.sidebar.markdown(f"**Q:** {q}\n\n**A:** {a}")
